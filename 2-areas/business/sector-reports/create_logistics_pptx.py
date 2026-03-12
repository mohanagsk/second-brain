#!/usr/bin/env python3
"""Generate Logistics & Supply Chain v2 presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets, emoji=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}" if emoji else title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.space_after = Pt(10)

def add_table_slide(prs, title, headers, rows, emoji=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}" if emoji else title
    p.font.size = Pt(28)
    p.font.bold = True
    
    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.3), Inches(0.9), Inches(9.4), Inches(0.4 * num_rows)).table
    
    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(11)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
add_title_slide(prs, "🚛 Logistics & Supply Chain", 
                "Business Idea Analyzer v2 | Feb 6, 2026\nShashank Framework Deep Dive")

# Slide 2: The Story
add_content_slide(prs, "The Story", [
    "India spends 14-16% of GDP on logistics (vs 8% in developed nations)",
    "$350B+ market with massive inefficiency — trucks run empty 40% of time",
    "National Logistics Policy aims to cut costs to 8% of GDP",
    "That's $100B+ in inefficiency waiting to be captured",
    "PM Gati Shakti digitally integrates 16 ministries for infra planning",
    "Budget 2026: ₹5.98 lakh crore allocated for transport sector"
], "📖")

# Slide 3: Market Numbers
add_table_slide(prs, "Market Numbers", 
    ["Segment", "2025", "2030", "CAGR"],
    [
        ["E-commerce Logistics", "$6.65B", "$10.29B", "9.11%"],
        ["D2C Logistics", "$7.55B", "$10.29B", "6.39%"],
        ["CEP (Courier Express)", "$5.68B", "$9.35B", "10.51%"],
        ["Total Logistics Market", "$350B+", "$500B+", "~8%"],
    ], "💰")

# Slide 4: Key Players & Gaps
add_table_slide(prs, "Competitive Landscape",
    ["Company", "Focus", "Strength", "Gap to Exploit"],
    [
        ["Delhivery", "Full-stack", "Scale, tech", "SMB neglect"],
        ["XpressBees", "E-commerce", "Speed", "Profitability"],
        ["Shiprocket", "Aggregation", "SMB easy", "Just matching"],
        ["Locus", "Route AI SaaS", "Tech moat", "Enterprise only"],
        ["Traditional 3PLs", "Physical ops", "Relationships", "No tech"],
    ], "🏆")

# Slide 5: Stakeholder Map
add_content_slide(prs, "Stakeholder Ecosystem", [
    "REGULATORS: MoRTH, GST Council, State RTOs — Control compliance",
    "SUPPLIERS: Truck OEMs, fuel companies, packaging — Input costs",
    "OPERATORS: Shippers, carriers, 3PLs, freight forwarders — Core ops",
    "CUSTOMERS: E-commerce, D2C, manufacturers, retailers — Demand",
    "ENABLERS: GPS/IoT providers, telecom, maps — Infrastructure",
    "FINANCE: Banks, NBFCs, factoring companies — Cash flow",
    "HIDDEN: Transport associations, CAs, insurance — Influence"
], "🗺️")

# Slide 6: People to Talk To
add_table_slide(prs, "Validation Interview Targets",
    ["Who", "Why", "Key Question"],
    [
        ["Fleet Owner (1-5 trucks)", "Lives fragmentation pain", "How do you find loads?"],
        ["Long-haul Driver", "Ground truth", "What's your daily struggle?"],
        ["SMB Manufacturer", "Ignored shipper", "How do you ship? What breaks?"],
        ["Ex-Delhivery Ops Manager", "Scale insights", "What's hardest to solve?"],
        ["Traditional Freight Broker", "Trust dynamics", "How do you match loads?"],
        ["Failed Logistics Founder", "Expensive lessons", "What killed your startup?"],
        ["E-commerce D2C Seller", "Multi-carrier user", "Why switch shipping partners?"],
    ], "🎤")

# Slide 7: GSRO - Government Data
add_content_slide(prs, "Government & Research Data", [
    "National Logistics Policy: Target <8% of GDP logistics cost",
    "PM Gati Shakti: 16 ministries integrated, infrastructure planning",
    "ULIP: 35+ logistics APIs — build on government data layer",
    "E-way Bill: 100M+ per month — compliance automation opportunity",
    "Sagarmala (ports) + Bharatmala (highways) — corridors improving",
    "Key sources: MoRTH, ULIP, GST Portal, IBEF, Invest India"
], "🏛️")

# Slide 8: Top 3 Opportunities
add_content_slide(prs, "Top 3 Startup Opportunities", [
    "1. SMB LOGISTICS OS: WhatsApp-first TMS for 100-1000 shipments/mo manufacturers",
    "   → $50-200/mo SaaS, vernacular, ULIP-integrated",
    "",
    "2. FLEET OWNER FINTECH: Instant payment + credit for 1-10 truck owners",
    "   → Data from ops → credit scoring, interest spread + SaaS",
    "",
    "3. COLD CHAIN VISIBILITY: IoT + real-time monitoring for pharma/F&B",
    "   → Per-shipment fee, high wastage = high WTP"
], "🎯")

# Slide 9: Action Plan
add_content_slide(prs, "Immediate Action Plan", [
    "THIS WEEK:",
    "  • Talk to 3 fleet owners at local transport nagar",
    "  • Interview 2 SMB manufacturers who ship goods",
    "  • Sign up for Shiprocket/Porter to understand UX",
    "",
    "THIS MONTH:",
    "  • Map all competitors and their pricing",
    "  • Find 1 ex-Delhivery/Rivigo insider for context",
    "  • Read ULIP API documentation (ulip.dpiit.gov.in)",
    "  • Test hypothesis with 10 potential customers"
], "✅")

# Save
output_path = "research/daily-sector/2026-02-06-logistics-supply-chain-v2.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
