#!/usr/bin/env python3
"""Generate Real Estate & PropTech sector analysis presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(26, 58, 95)  # Dark blue
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets, emoji=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RgbColor(26, 58, 95)
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}" if emoji else title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(50, 50, 50)
        p.space_after = Pt(12)

def add_numbers_slide(prs, title, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RgbColor(26, 58, 95)
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"💰 {title}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Metrics grid
    cols = 2
    rows = (len(data) + 1) // 2
    box_width = Inches(4.2)
    box_height = Inches(1.3)
    
    for i, (label, value) in enumerate(data):
        col = i % cols
        row = i // cols
        x = Inches(0.5) + col * (box_width + Inches(0.3))
        y = Inches(1.5) + row * (box_height + Inches(0.2))
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(240, 240, 245)
        box.line.fill.background()
        
        # Value
        val_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), box_width - Inches(0.4), Inches(0.6))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RgbColor(26, 58, 95)
        
        # Label
        lbl_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), box_width - Inches(0.4), Inches(0.4))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(100, 100, 100)

def add_comparison_slide(prs, title, winners, failures):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RgbColor(26, 58, 95)
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Winners column
    win_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.2), Inches(0.5))
    tf = win_header.text_frame
    p = tf.paragraphs[0]
    p.text = "🚀 WINNERS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(46, 139, 87)
    
    win_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.2), Inches(4.5))
    tf = win_box.text_frame
    for i, item in enumerate(winners):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # Failures column
    fail_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.2), Inches(0.5))
    tf = fail_header.text_frame
    p = tf.paragraphs[0]
    p.text = "💀 GRAVEYARD"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(178, 34, 34)
    
    fail_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.2), Inches(4.5))
    tf = fail_box.text_frame
    for i, item in enumerate(failures):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, "Real Estate & PropTech", "Business Idea Analyzer v2 | Feb 7, 2026")
    
    # Slide 2: Market Numbers
    add_numbers_slide(prs, "Market Numbers", [
        ("Market Size (2024)", "$570-584B"),
        ("Projected (2030)", "$1 Trillion"),
        ("PropTech Funding (2025)", "$550M"),
        ("CAGR", "7.3% - 16.6%"),
        ("Square Yards Revenue", "₹1,410 Cr"),
        ("NoBroker Users", "20M+"),
    ])
    
    # Slide 3: The Story
    add_content_slide(prs, "The Story", [
        "$570B industry still runs on WhatsApp forwards and paper documents",
        "Buying an apartment: 3 weeks to find + verify ownership + RERA + loan + registration",
        "Each step = different office, different process, different bribes",
        "The gap isn't software — it's TRUST INFRASTRUCTURE",
        "Whoever builds 'Razorpay of real estate' wins",
    ], "📖")
    
    # Slide 4: Winners vs Graveyard
    add_comparison_slide(prs, "🏆 Winners vs 💀 Graveyard",
        winners=[
            "NoBroker: 20M users, became a verb in Bangalore",
            "Square Yards: ₹1,410Cr revenue, full-stack, IPO planned",
            "Livspace: Captured post-purchase (interiors)",
            "PropertyShare: Leading SEBI-regulated fractional",
            "Pattern: Own one wedge deeply → expand to adjacent",
        ],
        failures=[
            "Grabhouse: Couldn't sustain no-brokerage, burned cash",
            "Commonfloor: Lost listings race to 99acres/Magicbricks",
            "2023 shutdowns: Over-raised in 2021, no path to profit",
            "Lesson: No-brokerage fails unless you replace broker VALUE",
            "Lesson: Listings marketplaces are winner-take-all",
        ]
    )
    
    # Slide 5: Deep Dive - Title Verification
    add_content_slide(prs, "Deep Dive: Title Verification", [
        "80% of property court disputes = title/ownership issues",
        "Land records scattered: 28 states, 28 different systems",
        "Sub-registrar + Revenue dept + RERA + Municipal = no single source",
        "DILRMP (Digital India Land Records) completes March 2026",
        "Banks pay ₹500-2000 per manual verification",
        "Automate 50% = ₹1,000 Cr+ addressable market",
        "Title Insurance barely exists in India (huge in US)",
    ], "🔬")
    
    # Slide 6: Stakeholder Map
    add_content_slide(prs, "Stakeholder Mapping", [
        "REGULATORS: RERA, State Revenue, SEBI (for REITs)",
        "DEVELOPERS: Builders, real estate companies (slow sales pain)",
        "PLATFORMS: NoBroker, 99acres, Square Yards",
        "BUYERS/TENANTS: Individuals and corporates",
        "ENABLERS: Banks, Lawyers, Brokers, Architects, Interior designers",
        "HIDDEN: NRIs (20% investment), Agricultural land, Industrial/warehouse",
        "KEY INSIGHT: Build for enablers (banks, lawyers) — they pay better than consumers",
    ], "🗺️")
    
    # Slide 7: People to Interview
    add_content_slide(prs, "10 People to Interview", [
        "First-time home buyer (0-6 months) — fresh pain memory",
        "Broker with 10+ years — knows every workaround",
        "Bank home loan manager — sees fraud patterns",
        "RERA official/consultant — regulatory inside view",
        "Failed PropTech founder — what went wrong",
        "Housing society secretary — collective decision pain",
        "Property lawyer (senior) — title dispute patterns",
        "NRI property investor — remote management pain",
        "Mid-size developer — sales and compliance pain",
        "Fractional platform user — new investor behavior",
    ], "🎤")
    
    # Slide 8: Government & Research
    add_content_slide(prs, "GSRO: Policy & Data", [
        "RERA (2016): Regulates developers, protects buyers — mandatory compliance",
        "DILRMP: All land records digital by March 2026 — KEY UNLOCK",
        "SM REIT (2024): Legitimizes fractional ownership platforms",
        "ULPIN: Unique Land Parcel ID (Aadhaar for land) — coming soon",
        "Data sources: RERA portals, IGRS, Knight Frank, JLL, CBRE reports",
        "March 2026 = potential 'Jio moment' for PropTech",
    ], "🏛️")
    
    # Slide 9: Top 3 Opportunities
    add_content_slide(prs, "Top 3 Solo Founder Opportunities", [
        "1. TITLE VERIFICATION API (B2B)",
        "   → Banks pay ₹500-2000/verification, start with Karnataka",
        "   → First customer: NBFC doing home loans",
        "",
        "2. NRI PROPERTY MANAGEMENT",
        "   → 20% of investment, zero good solutions",
        "   → Wedge: Dubai/US NRI WhatsApp groups",
        "",
        "3. RERA COMPLIANCE SaaS",
        "   → ₹50K-2L/project, developers hate manual compliance",
        "   → Start with one mid-size builder",
    ], "🎯")
    
    # Save
    output_path = os.path.expanduser("~/research/daily-sector/2026-02-07-real-estate-proptech-v2.pptx")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")

if __name__ == "__main__":
    main()
