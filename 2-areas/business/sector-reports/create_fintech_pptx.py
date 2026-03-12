#!/usr/bin/env python3
"""Generate FinTech & Payments Infrastructure presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(15, 23, 42)  # Slate 900
ACCENT = RGBColor(16, 185, 129)  # Emerald 500
WHITE = RGBColor(255, 255, 255)
LIGHT_TEXT = RGBColor(203, 213, 225)  # Slate 300

def add_dark_slide(title_text, content_items=None):
    """Add a dark-themed slide with title and bullet content."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Content
    if content_items:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()
            para.text = f"• {item}"
            para.font.size = Pt(24)
            para.font.color.rgb = LIGHT_TEXT
            para.space_after = Pt(12)
    
    return slide

# Slide 1: Title
slide1 = add_dark_slide("")
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "💰 FinTech & Payments Infrastructure"
p.font.size = Pt(52)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Business Idea Analyzer v2 | Feb 11, 2026"
p2.font.size = Pt(28)
p2.font.color.rgb = ACCENT
p2.alignment = PP_ALIGN.CENTER

# Slide 2: Market Numbers
add_dark_slide("📊 Market Numbers", [
    "UPI Transactions (2025): 228 billion",
    "UPI Value (2025): ₹300 trillion (~$3.6T)",
    "India Digital Payments (2025): $6.83B",
    "Projected (2034): $33.5B (16.1% CAGR)",
    "Embedded Payments Global: $3.5T transaction volume",
    "Fintech Unicorns in India: 22+"
])

# Slide 3: The Infrastructure Stack
add_dark_slide("🔬 Infrastructure Stack", [
    "Layer 1: Identity (Aadhaar, eKYC, Video KYC)",
    "Layer 2: Core Banking (TCS, Infosys — 15-20yr old systems)",
    "Layer 3: Payment Rails (UPI, IMPS, NEFT, RuPay, BBPS)",
    "Layer 4: Applications (PhonePe, Razorpay, Cashfree)",
    "",
    "AI Revolution: 60% lenders using ML for underwriting",
    "Our Angle: Build the 'boring' infrastructure everyone needs"
])

# Slide 4: Key Trends 2026
add_dark_slide("📰 Key Trends 2026", [
    "🔥 Credit Line on UPI (CLOU) — BNPL via UPI, banks scrambling",
    "🔥 Cross-Border Goes Live — Freelancer economy opportunity",
    "🔥 Authentication Overhaul — Beyond OTP, biometrics",
    "🔥 SRO Framework — Industry co-regulation with RBI",
    "",
    "Signal: Compliance is non-negotiable post-Paytm"
])

# Slide 5: Stakeholder Map
add_dark_slide("🗺️ Stakeholder Map", [
    "Regulators: RBI, NPCI, SEBI, IRDAI, MEITY",
    "Banks: SBI, HDFC (legacy pain) + SFBs (no tech talent)",
    "Fintechs: PhonePe, Razorpay, Cashfree (licensed)",
    "Enablers: AWS, TCS, Infosys, cloud providers",
    "Merchants: 63M+ SMBs, kiranas to enterprise",
    "Hidden: CAs (influence decisions), NPCI (makes rules)"
])

# Slide 6: People to Talk To
add_dark_slide("🎤 People to Talk To", [
    "1. SMB CFO — 'How many hrs/week on reconciliation?'",
    "2. NBFC Compliance Officer — 'What keeps you up at night?'",
    "3. Ex-Razorpay PM — 'What would you do differently?'",
    "4. Kirana Owner (Tier 2) — 'Why still prefer cash?'",
    "5. Bank Tech Head — 'Why is integration so slow?'",
    "6. RBI Consultant — 'What's coming in 2026-27?'"
])

# Slide 7: GSRO Data
add_dark_slide("🏛️ GSRO (Govt & Research)", [
    "Key RBI Regulations:",
    "  - PA/PG Guidelines, SRO Framework (2024)",
    "  - Authentication Directions (2025)",
    "  - Account Aggregator, DPDP Act",
    "",
    "Govt Initiatives: Digital India, PIDF, UPI International",
    "Data Sources: rbi.org.in, npci.org.in, pib.gov.in"
])

# Slide 8: Winners vs Graveyard
add_dark_slide("🚀 Winners vs 💀 Graveyard", [
    "WINNERS:",
    "  PhonePe (400M+ users), Razorpay ($7.5B), Zerodha (profitable Day 1)",
    "",
    "GRAVEYARD:",
    "  Paytm Bank (RBI restrictions) — Compliance > Growth",
    "  BharatPe (governance) — Co-founder conflicts kill",
    "  BNPL Players (defaults) — Easy credit = easy defaults"
])

# Slide 9: Top 3 Picks
add_dark_slide("🎯 Top 3 Picks for Divy", [
    "#1: Reconciliation SaaS for SMBs",
    "  Universal pain, clear WTP, existing tools suck",
    "",
    "#2: Compliance Copilot",
    "  Post-Paytm everyone scared of RBI. Reg tracker + alerts",
    "",
    "#3: WhatsApp Collections Bot",
    "  NBFCs bleeding on collections. Pay-for-success model"
])

# Save
output_path = "research/daily-sector/2026-02-11-fintech-deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
